from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ppt = Presentation("./file_load/fixtures/test_ppt.pptx")

for slide_number, slide in enumerate(ppt.slides, start=1):
    print(f"Slide {slide_number}:")
    for shape in slide.shapes:
        if shape.has_text_frame:  # 文本信息
            print(shape.text)
        
        if shape.has_table:  # 表格信息
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text
                    print(f"Row {row_idx + 1}, Column {col_idx + 1}: {cell_text}")
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE: # 图片信息
            imgae = shape.image
            image_filename = "./file_load/fixtures/pic_from_ppt.jpg"
            with open(image_filename, 'wb') as f:
                f.write(imgae.blob)